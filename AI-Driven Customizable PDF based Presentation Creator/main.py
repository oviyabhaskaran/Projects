
import os
import json
import time
import requests
from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from groq import Groq
from pptx import Presentation
from pptx.util import Pt  # for setting font size
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Function to call Groq LLM for summarization
def call_groq(prompt):
    url = "http://llama38b"  # Add your actual API URL here
    payload = json.dumps({
        "prompt": prompt,
        "max_tokens": 2048,
        "temperature": 0.1,
        "top_p": 0.95
    })
    headers = {'Content-Type': 'application/json'}
    response = requests.post(url, data=payload, headers=headers)

    if response.status_code == 200:
        response_json = response.json()
        print("Groq LLM response JSON:", response_json)  # Debug: Print the JSON response
        if isinstance(response_json, list):
            return response_json[0].get('text', 'No text output received')
        else:
            return response_json.get('text_output', 'No text output received')
    else:
        print(f"Error: {response.status_code} - {response.text}")  # Debug: Print error details
        return f"Error: {response.status_code} - {response.text}"


def clean_extracted_text(text):
    # Function to remove unwanted lines or characters
    cleaned_lines = []
    for line in text.splitlines():
        if line.startswith("N17/1/AYENG/HP1/ENG/TZ0/XX") or line.startswith("Choose either question") or line.strip().isdigit():
            continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines)

def summarize_pdf(file_path, slide_titles):
    if not os.path.isfile(file_path):
        raise ValueError(f"File path {file_path} is not a valid file or url")

    loader = PyPDFLoader(file_path)
    docs = loader.load()

    # Clean each page content in docs
    cleaned_docs = [clean_extracted_text(doc.page_content) for doc in docs]

    # Join the cleaned documents into one large text
    combined_text = "\n".join(cleaned_docs)
    print("length of combined text", len(combined_text))

    # Split the cleaned PDF content into smaller chunks using recursive chunking
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    start_time = time.perf_counter()
    texts = text_splitter.split_text(combined_text)
    print("no of texts ", len(texts))

    # Process in batches to avoid context length exceeded error
    intermediate_summaries = []
    for i in range(0, len(texts), 3):  # Adjust batch size if necessary
        batch_texts = texts[i:i+3]
        combined_batch_text = " ".join(batch_texts)

        prompt_template = """Based on the text below, generate content for the given slide titles in bullet points:
        "{text}"

        Slide Titles:
        {titles}

        FORMAT: {{'Slide Title': 'Slide Content'}}
        """
        slide_titles_str = "\n".join(slide_titles)
        prompt = PromptTemplate.from_template(prompt_template)
        prompt_str = prompt.format(text=combined_batch_text, titles=slide_titles_str)

        response = call_groq(prompt_str)
        #print("response is", response)
        summary_dict = parse_response_to_dict(response)
        #print("summary_dict", summary_dict)
        intermediate_summaries.append(summary_dict)
        #print("intermediate_summaries", intermediate_summaries)

    combined_summary = " ".join(" ".join(content.values()) for content in intermediate_summaries)
    #print("combined_summary are", combined_summary)

    # Create the prompt to include slide titles and contents
    final_prompt_template = """Based on the text below, generate content for the given slide titles in bullet points:
        "{text}"

        Slide Titles:
        {titles}

        FORMAT: {{'Slide Title': 'Slide Content'}}
    """
    final_prompt = PromptTemplate.from_template(final_prompt_template)
    print("final_prompt", final_prompt)
    final_prompt_str = final_prompt.format(text=combined_summary, titles=slide_titles_str)
    print("final_prompt_str", final_prompt_str)

    # Get the final response from the LLM
    final_response = call_groq(final_prompt_str)
    print("final_reponse is", final_response)

    end_time = time.perf_counter()
    Time_taken = end_time - start_time
    print("Total Time Taken", Time_taken)

    # Parse response to dictionary
    final_summary_dict = parse_response_to_dict(final_response)
    print("final_summary_dict", final_summary_dict)

    return final_summary_dict

def parse_response_to_dict(response):
    try:
        return json.loads(response)
    except json.JSONDecodeError:
        slides = {}
        lines = response.split("\n")
        current_slide = None
        current_content = []

        for line in lines:
            line = line.strip()
            if (line.startswith("**") and line.endswith("**")) or \
               (line.startswith("'") and line.endswith("'")) or \
               (line.startswith(""") and line.endswith(""")) or \
               (line.startswith("`") and line.endswith("`")):
                if current_slide:
                    slides[current_slide] = "\n".join(current_content)
                current_slide = line[2:-2] if line.startswith("**") else line[1:-1]
                current_content = []
            elif line:
                current_content.append(line)

        if current_slide:
            slides[current_slide] = "\n".join(current_content)

        return slides

# Function to generate main title and subtitle
def generate_main_title_and_subtitle(content):
    # Create the prompt to generate a main title and subtitle for the presentation
    prompt_template = """Based on the content below, generate a main title and subtitle for the presentation:
    "{text}"

    FORMAT: "**Main Title:** MAIN_TITLE\n**Subtitle:** SUBTITLE"
    """
    prompt = PromptTemplate.from_template(prompt_template)
    prompt_str = prompt.format(text=content)

    response = call_groq(prompt_str)
    print("Groq Response for Title and Subtitle:", response)  # Debug: Print response

    # Extract main title and subtitle from response
    lines = response.splitlines()
    main_title = None
    subtitle = None
    for line in lines:
        if line.startswith("**Main Title:** "):
            main_title = line[len("**Main Title:** "):].strip()
        elif line.startswith("**Subtitle:** "):
            subtitle = line[len("**Subtitle:** "):].strip()

    print("Extracted Main Title:", main_title)
    print("Extracted Subtitle:", subtitle)

    if main_title is None or main_title.strip() == '':
        raise ValueError("Main title could not be generated or is missing from response.")

    return main_title, subtitle
def remove_existing_text_boxes(slide):
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame and not shape == slide.shapes.title:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        slide.shapes._spTree.remove(shape._element)

# Function to insert response into PowerPoint template
def create_presentation(template_path, output_path, main_title, subtitle, slides_content, two_column_slides):
    prs = Presentation(template_path)

    # Modify the first slide with the title and subtitle
    first_slide = prs.slides[0]
    title = first_slide.shapes.title
    subtitle_shape = first_slide.placeholders[1]

    title.text = main_title
    subtitle_shape.text = subtitle

    # Set font style and size for main title
    title.text_frame.text = main_title
    title_text_format = title.text_frame.paragraphs[0].font
    title_text_format.name = 'Arial'  # font style
    title_text_format.size = Pt(32)  # font size
    title_text_format.bold = True
    title_text_format.color.rgb = RGBColor(0, 0, 0)  # font color (black)

    # Set font style and size for subtitle
    subtitle_shape.text_frame.text = subtitle
    subtitle_text_format = subtitle_shape.text_frame.paragraphs[0].font
    subtitle_text_format.name = 'Arial'  # font style
    subtitle_text_format.size = Pt(24)  # font size
    subtitle_text_format.color.rgb = RGBColor(100, 100, 100)  # font color (gray)

    # Add slides with titles and content from LLM response
    slide_layout = prs.slide_layouts[1]  # Using the title and content layout

    def add_slide(prs, title_text, content_text, is_two_column):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title

        # Remove existing text boxes except for the title
        remove_existing_text_boxes(slide)

        title.text = title_text

        # Set font style and size for slide title
        title_text_format = title.text_frame.paragraphs[0].font
        title_text_format.name = 'Arial'  # font style
        title_text_format.size = Pt(28)  # font size
        title_text_format.bold = True
        title_text_format.color.rgb = RGBColor(0, 0, 0)  # Set font color (black)

        if is_two_column:
            # Split content into two columns
            lines = content_text.split('\n')
            midpoint = len(lines) // 2
            left_column_content = "\n".join(lines[:midpoint])
            right_column_content = "\n".join(lines[midpoint:])

            # Create left text frame
            left_text_frame = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(4.5), Inches(5)
            ).text_frame
            left_text_frame.word_wrap = True

            for line in left_column_content.split('\n'):
                p = left_text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Setting the bullet point level

                # Set font style and size for slide content
                content_text_format = p.font
                content_text_format.name = 'Arial'  # font style
                content_text_format.size = Pt(18)  # font size
                content_text_format.color.rgb = RGBColor(50, 50, 50)  # font color (dark gray)

            # Create right text frame
            right_text_frame = slide.shapes.add_textbox(
                Inches(5.5), Inches(1.5), Inches(4.5), Inches(5)
            ).text_frame
            right_text_frame.word_wrap = True

            for line in right_column_content.split('\n'):
                p = right_text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Setting the bullet point level

                # Set font style and size for slide content
                content_text_format = p.font
                content_text_format.name = 'Arial'  # font style
                content_text_format.size = Pt(18)  # font size
                content_text_format.color.rgb = RGBColor(50, 50, 50)  # font color (dark gray)
        else:
            # Single column content
            text_frame = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(5)
            ).text_frame
            text_frame.word_wrap = True

            for line in content_text.split('\n'):
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Setting the bullet point level

                # Set font style and size for slide content
                content_text_format = p.font
                content_text_format.name = 'Arial'  # font style
                content_text_format.size = Pt(18)  # font size
                content_text_format.color.rgb = RGBColor(50, 50, 50)  # font color (dark gray)

    max_words_per_slide = 70  # Adjust this number based on your template's capacity

    for slide_title, slide_content in slides_content.items():
        is_two_column = slide_title in two_column_slides
        words = slide_content.split()
        if len(words) > max_words_per_slide:
            chunks = [" ".join(words[i:i + max_words_per_slide]) for i in range(0, len(words), max_words_per_slide)]
            for i, chunk in enumerate(chunks):
                if i == 0:
                    add_slide(prs, slide_title, chunk, is_two_column)
                else:
                    add_slide(prs, "To be contd..", chunk, is_two_column)
        else:
            add_slide(prs, slide_title, slide_content, is_two_column)

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def add_image_slide(prs, image_path, title):
    slide_layout = prs.slide_layouts[5]  # Choose an appropriate layout for the image slide (Title + Content)
    slide = prs.slides.add_slide(slide_layout)

    # Set slide title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set font style and size for slide title
    title_text_format = title_shape.text_frame.paragraphs[0].font
    title_text_format.name = 'Arial'  # font style
    title_text_format.size = Pt(28)  # font size
    title_text_format.bold = True
    title_text_format.color.rgb = RGBColor(0, 0, 0)  # font color (black)

    # Add image to the slide
    left_inch = Inches(1)
    top_inch = Inches(1.5)
    width_inch = Inches(8)
    height_inch = Inches(5)

    slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

def insert_image_slide(prs, image_path, title, position):
    slide_layout = prs.slide_layouts[5]  # Assuming slide layout 5 is a blank layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Set font style and size for slide title
    title_text_format = title_shape.text_frame.paragraphs[0].font
    title_text_format.name = 'Arial'
    title_text_format.size = Pt(28)
    title_text_format.bold = True
    title_text_format.color.rgb = RGBColor(0, 0, 0)  # Set font color (black)

    # Add image
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(image_path, left, top, width=Inches(8))

    # Move slide to the specified position
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide = slides[-1]
    xml_slides.remove(new_slide)
    xml_slides.insert(position - 1, new_slide)

# Function to load PowerPoint templates from a directory
def load_templates(directory):
    templates = {}
    for filename in os.listdir(directory):
        if filename.endswith(".pptx"):
            templates[filename] = os.path.join(directory, filename)
    return templates

# Function to select a PowerPoint template interactively
def select_template(templates):
    print("Available templates:")
    for i, template in enumerate(templates.keys(), 1):
        print(f"{i}. {template}")
    choice = int(input("Select a template by number: ")) - 1
    selected_template = list(templates.values())[choice]
    return selected_template

# Function to insert an image with its title into an existing slide
def insert_image_into_existing_slide(prs, image_path, image_title, slide_title, slides_content):
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text == slide_title:
            # Remove existing text boxes except for the title
            remove_existing_text_boxes(slide)

            title = slide.shapes.title
            title_text_format = title.text_frame.paragraphs[0].font
            title_text_format.name = 'Arial'
            title_text_format.size = Pt(28)
            title_text_format.bold = True
            title_text_format.color.rgb = RGBColor(0, 0, 0)  # Set font color (black)

            # Split the existing content and create a two-column layout
            content_text = slides_content[slide_title]
            left_text_frame = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(4.5), Inches(5)
            ).text_frame
            left_text_frame.word_wrap = True

            for line in content_text.split('\n'):
                p = left_text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Setting the bullet point level

                content_text_format = p.font
                content_text_format.name = 'Arial'
                content_text_format.size = Pt(18)
                content_text_format.color.rgb = RGBColor(50, 50, 50)  # font color (dark gray)

            # Add image and its title on the right side
            right_text_frame = slide.shapes.add_textbox(
                Inches(5.5), Inches(1.5), Inches(4.5), Inches(1)
            ).text_frame
            right_text_frame.word_wrap = True

            # Add image title
            p = right_text_frame.add_paragraph()
            p.text = image_title
            p.level = 0

            content_text_format = p.font
            content_text_format.name = 'Arial'
            content_text_format.size = Pt(18)
            content_text_format.bold = True
            content_text_format.color.rgb = RGBColor(0, 0, 0)  # font color (black)

            # Add image
            left = Inches(5.5)
            top = Inches(2.5)
            slide.shapes.add_picture(image_path, left, top, width=Inches(4.5))

def main():
    file_path = 'Knowledge_Graphs.pdf'

    print("Enter the slide titles for your presentation.")
    slide_titles = []
    while True:
        title = input("Enter a slide title (or type 'done' to finish): ")
        if title.lower() == 'done':
            break
        slide_titles.append(title)

    print("Reading and summarizing the PDF...")

    # User input for multiple image slide titles, image paths, and positions
    print("Getting images and placed in the specified slide position")
    image_titles_paths_positions = []
    while True:
        image_title = input("Enter the title for the image slide (or type 'done' to finish): ")
        if image_title.lower() == 'done':
            break
        image_path = input("Enter the path for the image: ")
        position = int(input("Enter the position (slide number) where the image slide should be inserted: "))
        image_titles_paths_positions.append((image_title, image_path, position))

    # Additional user input for images to be inserted into specific slides
    print("Getting image details and place alongside with the existed slide")
    images_for_existing_slides = []
    while True:
        image_title = input("Enter the title for the image (or type 'done' to finish): ")
        if image_title.lower() == 'done':
            break
        image_path = input("Enter the path for the image: ")
        slide_title = input("Enter the slide title where the image should be inserted: ")
        images_for_existing_slides.append((image_title, image_path, slide_title))

    llm_response_dict = summarize_pdf(file_path, slide_titles)

    # Combine all summarized content to generate the main title and subtitle
    combined_content = " ".join(llm_response_dict.values())
    try:
        main_title, subtitle = generate_main_title_and_subtitle(combined_content)
    except ValueError as e:
        print(f"Error: {e}")
        return

    # Get user input for which slides should be two columns
    print("Enter the titles of slides that should have two columns (type 'done' when finished):")
    two_column_slides = []
    while True:
        title = input("Enter a slide title which should be two split columns (or type 'done' to finish): ")
        if title.lower() == 'done':
            break
        two_column_slides.append(title.strip())

        # Directory containing PowerPoint templates
    templates_directory = r"templates"

    templates = load_templates(templates_directory)
    selected_template_path = select_template(templates)
    output_pptx_path = r"output_presentation.pptx"

    create_presentation(selected_template_path, output_pptx_path, main_title, subtitle, llm_response_dict, two_column_slides)

    # Load the created presentation
    prs = Presentation(output_pptx_path)

    # Add image slides to the presentation at specified positions
    for image_title, image_path, position in image_titles_paths_positions:
        insert_image_slide(prs, image_path, image_title, position)

    # Insert images into existing slides based on user input
    for image_title, image_path, slide_title in images_for_existing_slides:
        insert_image_into_existing_slide(prs, image_path, image_title, slide_title, llm_response_dict)

    # Save the modified presentation with the image slide
    prs.save(output_pptx_path)

    print(f"Presentation with image slide saved to {output_pptx_path}")

# Ensure the main function runs when the script is executed
if __name__ == "__main__":
    main()
